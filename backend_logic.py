import os
import json
import zipfile
import io
from pptx import Presentation
from google import genai
from google.genai import types

# --- 1. PPT 提取邏輯 (原 extract_content_v3) ---

def get_shape_text(shape):
    text = []
    if shape.has_text_frame:
        text.append(shape.text_frame.text)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text.append(cell.text_frame.text)
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            text.extend(get_shape_text(sub_shape).split('\n'))
    return "\n".join(filter(None, text))

def extract_text_from_ppt_stream(ppt_file):
    """Directly extracts structured text from the uploaded PPT file stream."""
    try:
        prs = Presentation(ppt_file)
        structured_output = []
        for index, slide in enumerate(prs.slides):
            slide_number = index + 1
            slide_content_parts = []
            for shape in slide.shapes:
                slide_content_parts.append(get_shape_text(shape))
            slide_content = "\n".join(filter(None, slide_content_parts))
            
            speaker_notes = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    speaker_notes = notes_text.strip()

            slide_block = f"""
============================================================
## SLIDE [{slide_number}] TEXT (Page {slide_number})
============================================================
{slide_content.strip()}

### SPEAKER NOTES
------------------------------------------------------------
{speaker_notes if speaker_notes else 'N/A (No Notes Found)'}
------------------------------------------------------------
"""
            structured_output.append(slide_block)
        return "\n".join(structured_output)
    except Exception as e:
        return f"Error extracting PPT: {str(e)}"

# --- 2. AI 分析邏輯 ---

def call_gemini_api(prompt, content, api_key, model_name='gemini-2.5-flash', output_json=False):
    """Generic function to call the Gemini API."""
    try:
        client = genai.Client(api_key=api_key)
        # Check if prompt is a full prompt template (e.g., mapping report) or needs content appended
        if "## MODULE CONTENT" in prompt or "## STRUCTURED PRESENTATION CONTENT" in prompt:
             full_prompt = prompt + "\n\n" + content
        else:
            full_prompt = content # For mapping, the JSONs are already included in the content string

        config = None
        if output_json:
            config = types.GenerateContentConfig(response_mime_type="application/json")

        response = client.models.generate_content(
            model=model_name,
            contents=full_prompt,
            config=config
        )
        # Ensure the response text doesn't contain code wrappers for Markdown output
        if not output_json:
            return response.text.strip().replace('```markdown', '').replace('```', '').strip()
            
        return response.text
    except Exception as e:
        return f"API Error: {str(e)}"

# --- PROMPTS 定義 (已優化，確保在後端邏輯中保持一致) ---

# Note: The mapping prompt uses string formatting, so we keep the placeholders.

PROMPT_SOURCE_ANALYSIS = """
## ROLE AND TASK
You are a Senior Technical Training Content Auditor. Extract a definitive list of all **new features, major enhancements, and significant UI/UX changes**.
## OUTPUT FORMAT
**STRICTLY** output a JSON list of objects.
## JSON STRUCTURE
Key: feature_id, feature_name_zh, feature_name_en, change_type, description (English), relevant_slides.
---
## STRUCTURED PRESENTATION CONTENT
"""

PROMPT_TARGET_INDEXING = """
## ROLE AND TASK
You are a Curriculum Architect. Map all key technical features and concepts presented in the module.
## OUTPUT FORMAT
**STRICTLY** output a JSON list of objects.
## JSON STRUCTURE
Key: target_concept_id, concept_name_zh, concept_name_en, concept_category, description_en, slide_range.
---
## MODULE CONTENT
"""

PROMPT_MAPPING_REPORT = """
## ROLE AND TASK
Compare a list of new product features (SOURCE) against an existing course's index (TARGET) to identify required updates.

## OUTPUT FORMAT
**STRICTLY** output a **Markdown** report.

## REPORT TEMPLATE (STRICTLY FOLLOW THIS)
# Curriculum Content Audit Report: [Course Name/Module Name]

## Course Module: [Full Course Name, e.g., ZCNE L1: Security Lv1 - Module 1-2]
## Audit Date: [Current Date, YYYY/MM/DD]
## Original Course Scope: Slides [Min Page Number] to [Max Page Number]
## Update Source: Nebula 19.20 New Feature List

---

## 總結摘要 (Summary)

[A concise 1-2 sentence summary of the overall impact of the new features on this specific module, in Chinese.]

---

## 📌 待辦事項清單 (To-Do List)

| 優先級 | 變更類型 | 待辦動作 | 相關功能 (EN/ZH) |
|---|---|---|---|
| HIGH/MEDIUM/LOW | [Change Type from Source List] | [Suggested Action] | [feature_name_en] / [feature_name_zh] |
| ... | ... | ... | ... |

---

## 📝 詳細稽核與映射報告 (Detailed Audit and Mapping)

### [N]. 【優先級變更】[Feature Name EN] ([Feature Name ZH])

* **更新來源 ID:** [feature_id]
* **來源描述:** [description]

| 舊課程概念 (Target Concept) | 概念 ID | 變更類型 | 建議更新內容 | 課程投影片範圍 |
|---|---|---|---|---|
| [concept_name_zh] | [target_concept_id] | [Change Type] | [Suggested Action based on mapping] | [slide_range] |
| [Related Concept 2...] | ... | ... | ... | ... |

[Continue for all relevant features...]

---
## SOURCE_LIST_JSON

{source_json}

## TARGET_INDEX_JSON

{target_json}
"""